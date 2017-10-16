import os
import re
import wget
import json
import random
import string
import shutil
import hashlib
import unicodedata
import datetime

from urllib import FancyURLopener
import urlparse

import urllib2
import simplejson

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

from giphypop import translate

from slide_weights import SlideWeights
from content_troll import Face

from pymarkovchain import MarkovChain
import Tkinter as tk

class MyOpener(FancyURLopener, object):
    version = "Mozilla/5.0 (Linux; U; Android 4.0.3; ko-kr; LG-L160L Build/IML74K) AppleWebkit/534.30 (KHTML, like Gecko) Version/4.0 Mobile Safari/534.30"


class Trollette:
    def __init__(self):
        self.presenter = ""
        self.title = ""

        self.slide_count = 0
        self.slide_min = 10
        self.slide_max = 12

        self.console = None
        self.output_dir = ""

        with open("terms.json", "r") as f:
            self.terms = json.load(f)

        with open(os.path.join("GIFs", "hashes.json"), "r") as f:
            self.gifs = json.load(f)

        with open(os.path.join("Images", "hashes.json"), "r") as f:
            self.images = json.load(f)

        # Load up the proverb data
        with open(os.path.join("Proverbs", "facts"), "r") as f:
            self.proverb_lines = f.readlines()
        self.proverbs = map(string.strip, self.proverb_lines)
        self.proverb_markov = MarkovChain("markov.db")
        self.proverb_markov.generateDatabase("".join(self.proverb_lines), n=1)

        # Make the text data
        # self.my_face = comptroller.face(self.title)
        # self.slide_titles = self.my_face.get_titles(50)
        # self.slide_bullets = self.my_face.get_bullets(100)

        self.my_face = Face("")

        self.slide_titles = ["shit", "balls", "butts"]
        self.slide_bullets = ["butts", "do", "stuff", "fucks", "more fucks"]

        self.ppt = Presentation()
        self.slide_weights = SlideWeights()

    def generate_slide_deck(self):
        # Create a place to put data and resources
        self.output_dir = os.path.join("Output", "%s_%s_%s" % (self.title,
                                                               self.presenter,
                                                               datetime.datetime.strftime(datetime.datetime.now(), '%Y_%m_%d_%H_%M_%S')))

        self.resources_dir = os.path.join(self.output_dir, "Resources")

        # Start with a fresh PowerPoint
        self.ppt = Presentation()

        # Make sure the directories exist
        try:
            os.makedirs(self.output_dir)
            os.makedirs(self.resources_dir)
        except:
            self.log("Directory %s already exists, overwriting..." % self.output_dir)

        self.slide_count = random.randint(self.slide_min, self.slide_max)
        self.log("Generating a slide deck of %d slides about %s" % (self.slide_count, self.title))

        try:
            self.log("Getting slide content...")
            self.my_face.set_topic(self.title)

            self.log("Generating slide titles...")
            self.slide_titles = self.my_face.get_titles(self.slide_count)

            self.log("Generating slide bullets...")
            self.slide_bullets = self.my_face.get_bullets(self.slide_count*3)
        except:
            self.log("Problem generating content for a talk on %s, exiting..." % self.title)
            return

        self.log_slide_weights()

        self.create_title_slide()
        self.create_slides()

        slide_path = os.path.join(self.output_dir, "%s.pptx" % self.title)
        self.ppt.save(slide_path)

        self.log("Successfully generated PPT on %s to %s" % (self.title, slide_path))

    def create_title_slide(self):
        title_slide_layout = self.ppt.slide_layouts[0]
        slide = self.ppt.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = self.title
        subtitle.text = self.presenter

    def create_slides(self):
        for i in range(self.slide_count):
            choice = self.slide_weights.choose_weighted()

            self.log("  Generating slide #%d: %s" % (i+1, choice))

            new_slide_layout = None
            if choice == "Single GIF":
                ns = self.create_gif_slide(random.choice(self.slide_titles), self.get_giphy_search_term(), i)
            elif choice == "Full Slide GIF":
                ns = self.create_full_gif_slide(self.get_giphy_search_term(), i)
            elif choice == "Single Image":
                ns = self.create_image_slide(random.choice(self.slide_titles), self.get_image_search_term(), i)
            elif choice == "Full Slide Image":
                ns = self.create_full_image_slide(self.get_image_search_term(), i)
            elif choice == "Information":
                ns = self.create_info_slide(i)
            elif choice == "Quotation":
                ns = self.create_quote_slide()

    def create_single_full_image_slide(self, image_path):
        blank_slide_layout = self.ppt.slide_layouts[6]
        new_slide = self.ppt.slides.add_slide(blank_slide_layout)

        left = Inches(0)
        top = Inches(0)
        height = Inches(8)
        width = Inches(10)
        pic = new_slide.shapes.add_picture(image_path, left, top, height=height, width=width)
        return new_slide

    def create_single_image_slide(self, slide_title, image_path):

        blank_slide_layout = self.ppt.slide_layouts[1]
        new_slide = self.ppt.slides.add_slide(blank_slide_layout)

        for shape in new_slide.shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format

                if phf.type == 1:
                    shape.text = slide_title

        left = Inches(1)
        top = Inches(1)
        height = Inches(6)
        width = Inches(8)
        pic = new_slide.shapes.add_picture(image_path, left, top, height=height, width=width)

        return new_slide

    def download_gif(self, term, slide_num):
        # If we have at least 3 local gifs, use one of those
        if (term in self.gifs) and (len(self.gifs[term]) > 3):
            return os.path.join("GIFs", "%s.gif" % random.choice(self.gifs[term]))

        try:
            # Download the gif
            img = translate(term)
            image_path = os.path.join(self.resources_dir, "%d.gif" % slide_num)
            wget.download(img.fixed_height.url, image_path)

            file_hasher = hashlib.md5()
            with open(image_path, "rb") as f:
                file_hasher.update(f.read())
            file_md5 = file_hasher.hexdigest()

            if not (term in self.gifs):
                self.gifs[term] = []

            if not (file_md5 in self.gifs[term]):
                self.gifs[term].append(file_md5)
                shutil.copy(image_path, os.path.join("GIFs", "%s.gif" % file_md5))
                with open(os.path.join("GIFs", "hashes.json"), "w") as f:
                    json.dump(self.gifs, f, indent=2)

            return image_path
        except:
            return None

    def download_image(self, term, slide_num):
        # If we have at least 3 local images, use one of those
        if (term in self.images) and (len(self.images[term]) > 3):
            return os.path.join("Images", "%s.img" % random.choice(self.images[term]))

        try:
            search_term = term
            if (random.randint(0, 100) % 2) == 0:
                search_term = self.title

            download_attempts = 0
            image_bytes = ""
            image_path = ""
            while download_attempts < 10:

                fetcher = urllib2.build_opener()
                start_index = random.randint(0, 50)
                search_url = "http://ajax.googleapis.com/ajax/services/search/images?v=1.0&q=%s&start=%s" % (search_term, str(start_index))
                f = fetcher.open(search_url)
                deserialized_output = simplejson.load(f)

                image_url = deserialized_output['responseData']['results'][random.randint(0, len(deserialized_output['responseData']['results'])-1)]['unescapedUrl']
                image_path = os.path.join(self.resources_dir, "%d.img" % slide_num)
                wget.download(image_url, image_path)

                with open(image_path, "rb") as f:
                    image_bytes = f.read()

                if (not image_bytes.startswith("<!DOCTYPE html>")) and (not image_bytes.startswith("<html>")):
                    break

                download_attempts += 1
                self.log("    Attempting to download image about %s failed try #%d" % (search_term, download_attempts))

            if image_bytes.startswith("<!DOCTYPE html") or image_bytes.startswith("<html>"):
                return None

            file_hasher = hashlib.md5()
            file_hasher.update(image_bytes)
            file_md5 = file_hasher.hexdigest()

            if not (term in self.images):
                self.images[term] = []

            if not (file_md5 in self.images[term]):
                self.images[term].append(file_md5)
                shutil.copy(image_path, os.path.join("Images", "%s.img" % file_md5))
                with open(os.path.join("Images", "hashes.json"), "w") as f:
                    json.dump(self.images, f, indent=2)

            return image_path
        except:
            return None

    def create_gif_slide(self, slide_title, term, slide_num):
        image_path = self.download_gif(term, slide_num)
        if image_path:
            return self.create_single_image_slide(slide_title, image_path)

    def create_full_gif_slide(self, term, slide_num):
        image_path = self.download_gif(term, slide_num)
        if image_path:
            return self.create_single_full_image_slide(image_path)

    def create_image_slide(self, slide_title, term, slide_num):
        while True:
            try:
                image_path = self.download_image(term, slide_num)
                if image_path:
                    return self.create_single_image_slide(slide_title, image_path)
            except:
                pass

    def create_full_image_slide(self, term, slide_num):
        image_path = self.download_image(term, slide_num)
        if image_path:
            return self.create_single_full_image_slide(image_path)

    def create_info_slide(self, slide_num):
        slide_title_info = random.choice(self.slide_titles)
        slide_title = slide_title_info
        if (random.randint(0, 100) % 3) == 0:
            slide_title = self.get_markov_proverb()

        sb = random.sample(self.slide_bullets, random.randint(1, 4))
        if (random.randint(0, 100) % 4) == 0:
            sb.append(self.get_markov_proverb())

        bullet_slide_layout = self.ppt.slide_layouts[1]
        new_slide = self.ppt.slides.add_slide(bullet_slide_layout)
        shapes = new_slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        body_shape.width = Inches(4)
        body_shape.left = Inches(1)
        body_shape.top = Inches(2)

        title_shape.text = slide_title

        tf = body_shape.text_frame
        for b in sb:
            p = tf.add_paragraph()
            #p.text = b

            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            run1 = p.add_run()
            run1.text = b
            font1 = run1.font
            font1.name = 'Sans Serif'
            font1.size = Pt(20)
            font1.italic = True
            font1.bold = True

        image_path = None
        attempts = 0
        while attempts < 10:
            try:
                tries = 0
                while (not image_path) and (tries < 10):
                    if (random.randint(0, 100) % 2) == 0:
                        search_term = self.get_giphy_search_term()
                        image_path = self.download_gif(search_term, slide_num)
                    else:
                        search_term = self.get_image_search_term()
                        image_path = self.download_image(search_term, slide_num)

                    tries += 1

                if tries < 10:
                    left = Inches(5.5)
                    top = Inches(3)
                    #height = Inches(3)
                    width = Inches(3)
                    pic = new_slide.shapes.add_picture(image_path, left, top, width=width)
                    break
                attempts += 1

            except:
                attempts += 1

        return new_slide

    def create_quote_slide(self):
        # Pick a random quote category and quote
        cat = random.choice(self.terms["quote_categories"])
        with open(os.path.join("Quotes", "quotes_%s.json" % cat)) as f:
            q1 = random.choice(json.load(f))

        cat = random.choice(self.terms["quote_categories"])
        with open(os.path.join("Quotes", "quotes_%s.json" % cat)) as f:
            q2 = random.choice(json.load(f))

        quote_text = "\"%s\"" % q1["quote"]
        if (random.randint(0,100) % 5) == 0:
            quote_text = random.choice(self.proverbs)

        quote_author = "- %s" % q2["name"]

        blank_slide_layout = self.ppt.slide_layouts[2]
        new_slide = self.ppt.slides.add_slide(blank_slide_layout)

        for shape in new_slide.shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.type == 1:
                    # Put in the quote title
                    shape.text = random.choice(self.terms["quote_titles"])

                elif phf.type == 2:
                    text_frame = shape.text_frame

                    # Create the quote text paragraph
                    p1 = text_frame.paragraphs[0]
                    p1.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    run1 = p1.add_run()
                    run1.text = quote_text
                    font1 = run1.font
                    font1.name = 'Sans Serif'
                    font1.size = Pt(30)
                    font1.italic = True
                    font1.bold = True

                    # Create the Author text paragraph
                    p2 = text_frame.add_paragraph()
                    p2.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
                    run2 = p2.add_run()
                    run2.text = quote_author
                    font2 = run2.font
                    font2.name = 'Calibri'
                    font2.size = Pt(24)

        return new_slide

    def get_giphy_search_term(self):
        st = random.choice(self.terms["giphy_searches"])
        if (random.randint(0, 100) % 5) == 0:
            st = self.title
        return st

    def get_image_search_term(self):
        st = random.choice(self.terms["image_searches"])
        if (random.randint(0, 100) % 2) == 0:
            st = self.title
        return st

    def get_proverb(self):
        return random.choice(self.proverb_lines)

    def get_markov_proverb(self, min=5, max=10):
        b = ""

        while True:
            b = self.proverb_markov.generateString()
            s = b.split(" ")
            if min <= len(s) <= max:
                break

        return b

    def add_term(self, term_type, term):
        if term in self.terms[term_type]:
            return "Term \"%s\" is already in %s!" % (term, term_type)
        else:
            self.terms[term_type].append(term)
            with open("terms.json", "w") as f:
                json.dump(self.terms, f, indent=4)
            return "Term \"%s\" added to %s." % (term, term_type)

    def delete_term(self, term_type, term):
        if not (term in self.terms[term_type]):
            return "Term \"%s\" isn't in %s, can't delete!" % (term, term_type)
        else:
            self.terms[term_type].remove(term)
            with open("terms.json", "w") as f:
                json.dump(self.terms, f, indent=4)
            return "Term \"%s\" removed from %s." % (term, term_type)

    def show_term_counts(self, term_type, term_json):
        log_str = "%s Terms:\n" % term_type
        for term in self.terms[term_type]:
            if term in term_json:
                log_str += "  %s: %d\n" % (term, len(term_json[term]))
            else:
                log_str += "  %s: 0\n" % term
        self.log(log_str)

    def get_file_md5(self, file_path):
        with open(file_path, "rb") as f:
            image_bytes = f.read()

        file_hasher = hashlib.md5()
        file_hasher.update(image_bytes)
        return file_hasher.hexdigest()

    def farm_image_term(self, term, amount=35, threshold=10):
        self.log("Farming images for %s..." % term)

        if not (term in self.images):
            self.images[term] = []

        attempt_count = 0
        while (attempt_count < threshold) and (len(self.images[term]) < amount):
            myopener = MyOpener()
            page = myopener.open('https://www.google.pt/search?q=%s&source=lnms&tbm=isch&sa=X&tbs=isz:l&tbm=isch' % term.replace(" ", "+"))
            html = page.read()

            for match in re.finditer(r'<a href="/imgres\?imgurl=(.*?)&amp;imgrefurl', html, re.IGNORECASE | re.DOTALL | re.MULTILINE):
                if len(self.images[term]) >= amount:
                    break

                try:
                    os.remove("test.img")
                except:
                    pass

                try:
                    path = urlparse.urlsplit(match.group(1)).path
                    self.log("  Downloading %s" % match.group(1))
                    myopener.retrieve(match.group(1), "test.img")

                    image_md5 = self.get_file_md5("test.img")

                    if not (image_md5 in self.images[term]):
                        self.images[term].append(image_md5)
                        shutil.copy("test.img", os.path.join("Images", "%s.img" % image_md5))
                        os.remove("test.img")
                        self.log("    Image saved to archive. %d/%d images." % (len(self.images[term]), amount))
                        attempt_count = 0
                    else:
                        self.log("    Already had image!")
                        attempt_count += 1
                except:
                    self.log("    Downloading failed")
                    attempt_count += 1

        self.log("Farming of %s images complete, now holding %d images" % (term, len(self.images[term])))

        with open(os.path.join("Images", "hashes.json"), "w") as f:
            json.dump(self.images, f, indent=2)

    def farm_images(self, amount=25, threshold=10):
        self.show_term_counts("image_searches", self.images)

        all_farm = self.terms["image_searches"]
        all_farm.extend(self.terms["talk_titles"])

        for term in all_farm:
            self.farm_image_term(term, amount, threshold)

    def farm_gif_term(self, term, amount=35, threshold=10):
        self.log("Farming GIFs for %s..." % term)

        if not (term in self.gifs):
            self.gifs[term] = []

        attempt_count = 0
        while (attempt_count < threshold) and (len(self.gifs[term]) < amount):

            image_path = "test.gif"
            try:
                os.remove(image_path)
            except:
                pass

            try:
                img = translate(term)
                wget.download(img.fixed_height.url, image_path)

                image_md5 = self.get_file_md5("test.gif")

                if not (image_md5 in self.gifs[term]):
                    self.gifs[term].append(image_md5)
                    shutil.copy(image_path, os.path.join("GIFs", "%s.gif" % image_md5))
                    self.log("    GIF saved to archive. %d/%d GIFs." % (len(self.gifs[term]), amount))
                    attempt_count = 0
                else:
                    self.log("    Already had GIF!")
                    attempt_count += 1
            except:
                self.log("    Downloading failed")
                attempt_count += 1

        self.log("Farming of %s GIFs complete, now holding %d GIFs" % (term, len(self.gifs[term])))

        with open(os.path.join("GIFs", "hashes.json"), "w") as f:
            json.dump(self.gifs, f, indent=2)

    def farm_gifs(self, amount=35, threshold=10):
        self.show_term_counts("giphy_searches", self.gifs)

        all_farm = self.terms["giphy_searches"]
        all_farm.extend(self.terms["talk_titles"])

        for term in all_farm:

            self.log("Farming GIFs for %s..." % term)

            if not (term in self.gifs):
                self.gifs[term] = []

            self.farm_gif_term(term, amount, threshold)

    def farm_content(self, all_content):
        for talk_title in self.terms["talk_titles"]:
            talk_path = os.path.join("Content", "%s.txt" % talk_title)
            # Either we're replacing all content or we're only replacing files that don't exist
            if all_content or (not os.path.exists(talk_path)):
                self.log("Farming data on %s..." % talk_title)
                with open(talk_path, "w") as f:
                    content = self.my_face.fully_research_topic(talk_title, self.log)
                    if type(content) is str:
                        clean_content = content
                    else:
                        clean_content = unicodedata.normalize('NFKD', content).encode('ascii', 'ignore')
                    f.write(clean_content)

    def log_slide_weights(self):
        self.log(self.slide_weights.get_weights_string())

    def log(self, message):
        if self.console:
            self.console.config(state=tk.NORMAL)
            self.console.insert(tk.END, "%s\n" % message)
            self.console.see(tk.END)
            self.console.config(state=tk.DISABLED)
            self.console.update()
        else:
            print(message)




def main():
    app = TrolletteGUI(tk.Tk())
    app.run()

if __name__ == "__main__":
    main()

__author__ = 'wartortell'