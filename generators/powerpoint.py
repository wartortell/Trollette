import os
import sys
import json
import random
import shutil
import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pymarkovchain import MarkovChain

from slide_weights import SlideWeights
from helpers import SlideWeights, get_data_folder


class PowerpointGenerator:
    def __init__(self, logger=None):
        
        self.logger = logger
        if not self.logger:
            self.logger = sys.stdout.write

        self.slide_weights = SlideWeights()

        with open("terms.json", "r") as f:
            self.terms = json.load(f)

        with open(os.path.join(get_data_folder("images"), "hashes.json"), "r") as f:
            self.images = json.load(f)

        with open(os.path.join(get_data_folder("gifs"), "hashes.json"), "r") as f:
            self.gifs = json.load(f)

        # Load up the proverb data
        with open(os.path.join("Proverbs", "facts"), "r") as f:
            self.proverb_lines = f.readlines()
        self.proverbs = self.proverb_lines
        self.proverb_markov = MarkovChain("markov.db")
        self.proverb_markov.generateDatabase("".join(self.proverb_lines), n=1)

    def generate_slide_deck(self, title, presenter, slide_min, slide_max):
        # Create a place to put data and resources
        self.output_dir = os.path.join("Output", "%s_%s_%s" % (title, presenter,
                                                               datetime.datetime.strftime(datetime.datetime.now(), '%Y_%m_%d_%H_%M_%S')))

        self.resources_dir = os.path.join(self.output_dir, "Resources")

        data_sources = self.terms["full_talk_titles"][title]
        data_jsons = []
        for d in data_sources:
            with open(os.path.join("Data", "{}.json".format(d)), "r") as f:
                data_jsons.append(json.load(f))

        # Start with a fresh PowerPoint
        self.ppt = Presentation()

        # Make sure the directories exist
        try:
            os.makedirs(self.output_dir)
            os.makedirs(self.resources_dir)
        except:
            self.logger("Directory %s already exists, overwriting..." % self.output_dir)

        slide_count = random.randint(slide_min, slide_max)
        self.logger("Generating a slide deck of %d slides about %s" % (slide_count, title))

        try:
            self.logger("Getting slide content...")
            #self.my_face.set_topic(self.title)

            self.logger("Generating slide titles...")
            slide_titles = []
            slide_bullets = []
            for j in data_jsons:
                if "titles" in j:
                    slide_titles.extend(j["titles"])
                if "bullets" in j:
                    slide_bullets.extend(j["bullets"])
        except Exception as e:
            self.logger("Problem generating content for a talk on %s, exiting..." % title)
            self.logger(str(e))
            return

        self.log_slide_weights()

        self.create_title_slide(title, presenter)
        self.create_slides(title, slide_titles, slide_bullets, random.randint(slide_min, slide_max))
        self.create_questions_slide()

        slide_path = os.path.join(self.output_dir, "%s.pptx" % title)
        self.ppt.save(slide_path)

        self.logger("Successfully generated PPT on %s to %s" % (title, slide_path))

        return slide_path

    def launch_deck(self, path):
        pass

    def create_title_slide(self, topic, presenter):
        self.logger("Creating Title Slide...")
        title_slide_layout = self.ppt.slide_layouts[0]
        slide = self.ppt.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = topic
        subtitle.text = presenter

    def create_bio_slide(self, presenter):
        if presenter:
            pass
        else:
            pass

    def create_slides(self, title, titles, bullets, count):

        last_choice = ""
        for i in range(count):
            choice = self.slide_weights.choose_weighted()

            if (last_choice == choice) and (choice != "Information"):
                i -= 1
                continue

            self.logger("  Generating slide #%d: %s" % (i+1, choice))

            if choice == "Single GIF":
                ns = self.create_gif_slide(random.choice(titles), self.get_giphy_search_term(title), i)
            elif choice == "Full Slide GIF":
                ns = self.create_full_gif_slide(self.get_giphy_search_term(title), i)
            elif choice == "Single Image":
                ns = self.create_image_slide(random.choice(titles), self.get_image_search_term(title), i)
            elif choice == "Full Slide Image":
                ns = self.create_full_image_slide(self.get_image_search_term(title), i)
            elif choice == "Information":
                ns = self.create_info_slide(i, title, titles, bullets)
            elif choice == "Tweet":
                ns = self.create_tweet_slide(title)
            elif choice == "Quotation":
                ns = self.create_quote_slide()
            elif choice == "Definition":
                ns = self.create_quote_slide()
            elif choice == "Urban Definition":
                ns = self.create_urban_slide(title)

            last_choice = choice

    def create_single_full_image_slide(self, image_path):
        blank_slide_layout = self.ppt.slide_layouts[6]
        new_slide = self.ppt.slides.add_slide(blank_slide_layout)

        left = Inches(0)
        top = Inches(0)
        height = Inches(8)
        width = Inches(10)
        pic = new_slide.shapes.add_picture(image_path, left, top, height=height, width=width)
        return new_slide

    def create_single_image_slide(self, slide_title, image_path):

        blank_slide_layout = self.ppt.slide_layouts[1]
        new_slide = self.ppt.slides.add_slide(blank_slide_layout)

        for shape in new_slide.shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format

                if phf.type == 1:
                    shape.text = slide_title

        left = Inches(1)
        top = Inches(2)
        height = Inches(5)
        width = Inches(8)
        pic = new_slide.shapes.add_picture(image_path, left, top, height=height, width=width)

        return new_slide

    def create_questions_slide(self):
        term = self.get_image_search_term("o shit waddup!!!!")
        image_path = os.path.join("Images", term, "{}.img".format(random.choice(self.images[term])))
        self.create_single_image_slide("Questions?", image_path)

    def create_gif_slide(self, slide_title, term, slide_num):
        image_path = os.path.join("Gifs", term, "{}.gif".format(random.choice(self.gifs[term])))
        if image_path:
            return self.create_single_image_slide(slide_title, image_path)

    def create_full_gif_slide(self, term, slide_num):
        image_path = os.path.join("Gifs", term, "{}.gif".format(random.choice(self.gifs[term])))
        if image_path:
            return self.create_single_full_image_slide(image_path)

    def create_image_slide(self, slide_title, term, slide_num):
        while True:
            try:
                image_path = os.path.join("Images", term, "{}.img".format(random.choice(self.images[term])))
                if image_path:
                    return self.create_single_image_slide(slide_title, image_path)
            except:
                pass

    def create_full_image_slide(self, term, slide_num):
        image_path = os.path.join("Images", term, "{}.img".format(random.choice(self.images[term])))
        if image_path:
            return self.create_single_full_image_slide(image_path)

    def create_info_slide(self, slide_num, title, titles, bullets):
        slide_title_info = random.choice(titles)
        slide_title = slide_title_info
        if (random.randint(0, 100) % 3) == 0:
            slide_title = self.get_markov_proverb()

        sb = random.sample(bullets, random.randint(1, 4))
        if (random.randint(0, 100) % 4) == 0:
            sb.append(self.get_markov_proverb())

        bullet_slide_layout = self.ppt.slide_layouts[1]
        new_slide = self.ppt.slides.add_slide(bullet_slide_layout)
        shapes = new_slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        body_shape.width = Inches(4)
        body_shape.left = Inches(1)
        body_shape.top = Inches(2)

        title_shape.text = slide_title

        tf = body_shape.text_frame
        for b in sb:
            p = tf.add_paragraph()
            #p.text = b

            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            run1 = p.add_run()
            run1.text = b
            font1 = run1.font
            font1.name = 'Sans Serif'
            font1.size = Pt(20)
            font1.italic = True
            font1.bold = True

        image_path = None
        attempts = 0
        while attempts < 10:
            try:
                tries = 0
                while (not image_path) and (tries < 10):
                    if (random.randint(0, 100) % 2) == 0:
                        search_term = self.get_giphy_search_term(title)
                        image_path = os.path.join("Gifs", search_term, "{}.gif".format(random.choice(self.gifs[search_term])))
                    else:
                        search_term = self.get_image_search_term(title)
                        image_path = os.path.join("Images", search_term, "{}.img".format(random.choice(self.images[search_term])))

                    tries += 1

                if tries < 10:
                    left = Inches(5.5)
                    top = Inches(3)
                    #height = Inches(3)
                    width = Inches(3)
                    pic = new_slide.shapes.add_picture(image_path, left, top, width=width)
                    break
                attempts += 1

            except:
                attempts += 1

        return new_slide

    def create_quote_slide(self):
        # Pick a random quote category and quote
        cat = random.choice(self.terms["quote_categories"])
        with open(os.path.join("Quotes", "quotes_%s.json" % cat)) as f:
            q1 = random.choice(json.load(f))

        cat = random.choice(self.terms["quote_categories"])
        with open(os.path.join("Quotes", "quotes_%s.json" % cat)) as f:
            q2 = random.choice(json.load(f))

        quote_text = "\"%s\"" % q1["quote"]
        if (random.randint(0,100) % 5) == 0:
            quote_text = random.choice(self.proverbs)

        quote_author = "- %s" % q2["name"]

        blank_slide_layout = self.ppt.slide_layouts[2]
        new_slide = self.ppt.slides.add_slide(blank_slide_layout)

        for shape in new_slide.shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.type == 1:
                    # Put in the quote title
                    shape.text = random.choice(self.terms["quote_titles"])

                elif phf.type == 2:
                    text_frame = shape.text_frame

                    # Create the quote text paragraph
                    p1 = text_frame.paragraphs[0]
                    p1.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    run1 = p1.add_run()
                    run1.text = quote_text
                    font1 = run1.font
                    font1.name = 'Sans Serif'
                    font1.size = Pt(30)
                    font1.italic = True
                    font1.bold = True

                    # Create the Author text paragraph
                    p2 = text_frame.add_paragraph()
                    p2.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
                    run2 = p2.add_run()
                    run2.text = quote_author
                    font2 = run2.font
                    font2.name = 'Calibri'
                    font2.size = Pt(24)

        return new_slide

    def create_tweet_slide(self, title):
        # Pick a random quote category and quote
        tweets = list()
        for query in self.terms["full_talk_titles"][title]:
            with open(os.path.join("Data", "{}.json".format(query)), "r") as f:
                query_json = json.load(f)
            tweets.extend(query_json["tweets"])

        tweet = random.choice(tweets)

        quote_text = "\"%s\"" % tweet["tweet"]
        quote_author = "- %s" % tweet["author"]

        blank_slide_layout = self.ppt.slide_layouts[2]
        new_slide = self.ppt.slides.add_slide(blank_slide_layout)

        for shape in new_slide.shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.type == 1:
                    # Put in the quote title
                    shape.text = random.choice(self.terms["tweet_titles"])

                elif phf.type == 2:
                    text_frame = shape.text_frame

                    # Create the quote text paragraph
                    p1 = text_frame.paragraphs[0]
                    p1.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    run1 = p1.add_run()
                    run1.text = quote_text
                    font1 = run1.font
                    font1.name = 'Sans Serif'
                    font1.size = Pt(30)
                    font1.italic = True
                    font1.bold = True

                    # Create the Author text paragraph
                    p2 = text_frame.add_paragraph()
                    p2.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
                    run2 = p2.add_run()
                    run2.text = quote_author
                    font2 = run2.font
                    font2.name = 'Calibri'
                    font2.size = Pt(24)

        return new_slide

    def create_urban_slide(self, title):
        # Pick a random quote category and quote
        urban_defs = list()
        for query in self.terms["full_talk_titles"][title]:
            with open(os.path.join("Data", "{}.json".format(query)), "r") as f:
                query_json = json.load(f)
            urban_defs.extend(query_json["defs"])

        urban_def = random.choice(urban_defs)

        quote_text = "\"%s\"" % urban_def

        blank_slide_layout = self.ppt.slide_layouts[2]
        new_slide = self.ppt.slides.add_slide(blank_slide_layout)

        for shape in new_slide.shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.type == 1:
                    # Put in the quote title
                    shape.text = random.choice(self.terms["full_talk_titles"][title])

                elif phf.type == 2:
                    text_frame = shape.text_frame

                    # Create the quote text paragraph
                    p1 = text_frame.paragraphs[0]
                    p1.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    run1 = p1.add_run()
                    run1.text = quote_text
                    font1 = run1.font
                    font1.name = 'Sans Serif'
                    font1.size = Pt(30)
                    font1.italic = True
                    font1.bold = True

        return new_slide

    def get_giphy_search_term(self, title):
        # Find a giphy search term that has downloaded gifs
        while True:
            st = random.choice(self.terms["giphy_searches"])
            if (random.randint(0, 100) % 5) == 0:
                st = random.choice(self.terms["full_talk_titles"][title])

            gif_dir = os.path.join("Gifs", st)
            if os.path.isdir(gif_dir) and len(os.listdir(gif_dir)):
                return st

    def get_image_search_term(self, title):
        # Find an image search term that has downloaded gifs
        while True:
            st = random.choice(self.terms["image_searches"])
            if (random.randint(0, 100) % 2) == 0:
                st = random.choice(self.terms["full_talk_titles"][title])

            image_dir = os.path.join("Images", st)
            if os.path.isdir(image_dir) and len(os.listdir(image_dir)):
                return st

    def get_proverb(self):
        return random.choice(self.proverb_lines)

    def get_markov_proverb(self, min=5, max=10):
        b = ""

        while True:
            b = self.proverb_markov.generateString()
            s = b.split(" ")
            if min <= len(s) <= max:
                break

        return b

    def log_slide_weights(self):
        self.logger(self.slide_weights.get_weights_string())


if __name__ == "__main__":
    df = PowerpointGenerator()
    df.generate_slide_deck("APTs Are Evil Hitler Nazi Robots", "Dick Cheese", 12, 15)
    #print(data)